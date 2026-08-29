# -*- coding: utf-8 -*-
"""Monta o deck de validacao do DisperSim 3D.

Requer python-pptx. Roda como:

    python docs/build_deck.py

e reescreve docs/DisperSim3D-Validacao.pptx. Os numeros vem de
docs/validation.md e docs/benchmark-results.md; ao atualizar um deles,
atualize o slide correspondente aqui.
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                   "DisperSim3D-Validacao.pptx")

# ── paleta ──────────────────────────────────────────────────────────
PAPER   = RGBColor(0xEC, 0xEF, 0xEE)
SURFACE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL   = RGBColor(0xE3, 0xE8, 0xE7)
INK     = RGBColor(0x11, 0x1A, 0x1D)
MUTED   = RGBColor(0x5B, 0x6A, 0x6E)
FAINT   = RGBColor(0x84, 0x96, 0x99)
RULE    = RGBColor(0xCD, 0xD6, 0xD4)
ACCENT  = RGBColor(0x14, 0x5E, 0x5A)
SIGNAL  = RGBColor(0xC2, 0x5B, 0x00)
GOOD    = RGBColor(0x2C, 0x7A, 0x5A)
WARN    = RGBColor(0x9A, 0x6A, 0x0C)
BAD     = RGBColor(0xA9, 0x32, 0x26)
TRACK   = RGBColor(0xD6, 0xDD, 0xDB)

F_TITLE = "Bahnschrift"
F_BODY  = "Segoe UI"
F_MONO  = "Consolas"

W, H = 13.333, 7.5

prs = Presentation()
prs.slide_width = Inches(W)
prs.slide_height = Inches(H)
BLANK = prs.slide_layouts[6]


# ── utilitarios ─────────────────────────────────────────────────────
def slide(bg=PAPER):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = bg
    r.line.fill.background(); r.shadow.inherit = False
    return s


def box(s, x, y, w, h, fill=None, line=None, lw=0.75):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        r.fill.background()
    else:
        r.fill.solid(); r.fill.fore_color.rgb = fill
    if line is None:
        r.line.fill.background()
    else:
        r.line.color.rgb = line; r.line.width = Pt(lw)
    r.shadow.inherit = False
    return r


def text(s, x, y, w, h, runs, size=14, font=F_BODY, color=INK, bold=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.15, space_after=0,
         caps=False, tracking=0):
    """runs: str, ou lista de paragrafos; cada paragrafo str ou lista de (txt, dict)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0

    paras = runs if isinstance(runs, list) else [runs]
    for i, p in enumerate(paras):
        para = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        para.alignment = align
        para.line_spacing = spacing
        para.space_after = Pt(space_after)
        pieces = p if isinstance(p, list) else [(p, {})]
        for txt, ov in pieces:
            run = para.add_run()
            run.text = txt.upper() if (ov.get("caps", caps)) else txt
            f = run.font
            f.name = ov.get("font", font)
            f.size = Pt(ov.get("size", size))
            f.bold = ov.get("bold", bold)
            f.color.rgb = ov.get("color", color)
            tr = ov.get("tracking", tracking)
            if tr:
                f._rPr.set("spc", str(int(tr * 100)))
    return tb


def header(s, eyebrow, title, sub=None):
    text(s, 0.85, 0.55, 11.6, 0.3, eyebrow, size=10.5, font=F_MONO,
         color=ACCENT, caps=True, tracking=1.4)
    text(s, 0.85, 0.92, 11.6, 0.75, title, size=30, font=F_TITLE, bold=True, color=INK)
    y = 1.78
    if sub:
        text(s, 0.85, y, 11.0, 0.5, sub, size=13, color=MUTED, spacing=1.3)
        y += 0.62
    box(s, 0.85, y, 11.63, 0.012, fill=RULE)
    return y + 0.34


def gauge(s, x, y, w, pos=None, span=None, color=GOOD):
    """Trilho log de 0,5x a 2,0x. pos/span em fracao 0..1."""
    box(s, x, y, w, 0.16, fill=TRACK)
    cx = x + w / 2.0
    box(s, cx - 0.008, y - 0.05, 0.016, 0.26, fill=SIGNAL)
    if span:
        a, b = sorted(span)
        box(s, x + w * a, y + 0.02, max(w * (b - a), 0.03), 0.12, fill=color)
    elif pos is not None:
        box(s, x + w * pos - 0.025, y + 0.02, 0.05, 0.12, fill=color)


def gauge_block(s, x, y, w, rows, label_w=3.5, obs_w=1.7, track_w=3.4):
    """rows: (nome, detalhe, medido, pos|span, valor, cor)"""
    text(s, x, y, label_w, 0.2, "Ensaio", size=8.5, font=F_MONO, color=FAINT, caps=True, tracking=1.0)
    text(s, x + label_w, y, obs_w, 0.2, "Medido", size=8.5, font=F_MONO, color=FAINT, caps=True, tracking=1.0)
    text(s, x + label_w + obs_w, y, track_w, 0.2, "0,5\u00d7   \u2190  1,0  \u2192   2,0\u00d7",
         size=8.5, font=F_MONO, color=FAINT, tracking=1.0)
    text(s, x + label_w + obs_w + track_w + 0.15, y, 0.85, 0.2, "Raz\u00e3o",
         size=8.5, font=F_MONO, color=FAINT, caps=True, tracking=1.0, align=PP_ALIGN.RIGHT)
    yy = y + 0.28
    box(s, x, yy, label_w + obs_w + track_w + 1.0, 0.012, fill=RULE)
    yy += 0.14
    for nome, det, obs, mark, val, cor in rows:
        text(s, x, yy, label_w - 0.15, 0.22, nome, size=11.5, color=INK)
        if det:
            text(s, x, yy + 0.21, label_w - 0.15, 0.2, det, size=8.5, font=F_MONO, color=FAINT)
        text(s, x + label_w, yy + 0.03, obs_w - 0.15, 0.22, obs, size=9.5, font=F_MONO, color=MUTED)
        if isinstance(mark, tuple):
            gauge(s, x + label_w + obs_w, yy + 0.07, track_w - 0.2, span=mark, color=cor)
        else:
            gauge(s, x + label_w + obs_w, yy + 0.07, track_w - 0.2, pos=mark, color=cor)
        text(s, x + label_w + obs_w + track_w + 0.05, yy + 0.03, 0.95, 0.22, val,
             size=11, font=F_MONO, bold=True, color=cor, align=PP_ALIGN.RIGHT)
        yy += 0.58 if det else 0.5
    return yy


def table(s, x, y, w, cols, rows, widths=None, fsize=10.5, rowh=0.32,
          aligns=None, headh=0.34):
    n_r, n_c = len(rows) + 1, len(cols)
    shp = s.shapes.add_table(n_r, n_c, Inches(x), Inches(y), Inches(w), Inches(headh + rowh * len(rows)))
    tbl = shp.table
    tbl.first_row = False
    tbl.horz_banding = False
    # remove o estilo padrao (azul)
    tblPr = tbl._tbl.find(qn('a:tblPr'))
    for e in tblPr.findall(qn('a:tableStyleId')):
        tblPr.remove(e)

    if widths:
        tot = sum(widths)
        for i, cw in enumerate(widths):
            tbl.columns[i].width = Emu(int(Inches(w) * cw / tot))
    tbl.rows[0].height = Inches(headh)
    for i in range(len(rows)):
        tbl.rows[i + 1].height = Inches(rowh)

    aligns = aligns or [PP_ALIGN.LEFT] * n_c

    def cell(rr, cc, txt, size, color, bold, font, fill, align, border_bottom):
        c = tbl.cell(rr, cc)
        c.margin_left = Inches(0.09); c.margin_right = Inches(0.09)
        c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
        c.vertical_anchor = MSO_ANCHOR.MIDDLE
        c.fill.solid(); c.fill.fore_color.rgb = fill
        tf = c.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; p.alignment = align
        r = p.add_run(); r.text = txt
        r.font.size = Pt(size); r.font.color.rgb = color
        r.font.bold = bold; r.font.name = font
        tcPr = c._tc.get_or_add_tcPr()
        for tag in ('a:lnL', 'a:lnR', 'a:lnT', 'a:lnB'):
            for e in tcPr.findall(qn(tag)):
                tcPr.remove(e)
        ln = tcPr.makeelement(qn('a:lnB'), {'w': '9525', 'cap': 'flat',
                                            'cmpd': 'sng', 'algn': 'ctr'})
        fillel = ln.makeelement(qn('a:solidFill'), {})
        clr = ln.makeelement(qn('a:srgbClr'), {'val': str(border_bottom)})
        fillel.append(clr); ln.append(fillel); tcPr.append(ln)

    for j, h in enumerate(cols):
        cell(0, j, h, fsize - 1.5, FAINT, False, F_MONO, PANEL, aligns[j], 'CDD6D4')
    for i, row in enumerate(rows):
        for j, v in enumerate(row):
            txt, ov = (v, {}) if isinstance(v, str) else v
            cell(i + 1, j, txt,
                 ov.get("size", fsize), ov.get("color", INK), ov.get("bold", False),
                 ov.get("font", F_BODY), ov.get("fill", SURFACE),
                 ov.get("align", aligns[j]), 'DDE3E1')
    return shp


def callout(s, x, y, w, h, eyebrow, paras, bar=ACCENT, size=12.5):
    box(s, x, y, w, h, fill=SURFACE, line=RULE)
    box(s, x, y, 0.045, h, fill=bar)
    text(s, x + 0.32, y + 0.22, w - 0.65, 0.25, eyebrow, size=10, font=F_MONO,
         color=bar, caps=True, tracking=1.4)
    text(s, x + 0.32, y + 0.56, w - 0.65, h - 0.76, paras, size=size, color=INK,
         spacing=1.32, space_after=8)


def footer(s, n):
    text(s, 0.85, 6.92, 6.0, 0.25, "DisperSim 3D  \u00b7  Doss\u00ea de valida\u00e7\u00e3o",
         size=8.5, font=F_MONO, color=FAINT, caps=True, tracking=1.2)
    text(s, 10.5, 6.92, 1.98, 0.25, str(n), size=8.5, font=F_MONO, color=FAINT,
         align=PP_ALIGN.RIGHT, tracking=1.2)


# ════════════════════════════════════════════════════════════════════
# 1. Capa
# ════════════════════════════════════════════════════════════════════
s = slide(INK)
box(s, 0, 0, W, 0.055, fill=SIGNAL)
text(s, 0.85, 1.0, 11.6, 0.3, "DisperSim 3D  \u00b7  Doss\u00ea de valida\u00e7\u00e3o",
     size=11, font=F_MONO, color=RGBColor(0x58, 0xBD, 0xB1), caps=True, tracking=1.6)
text(s, 0.85, 1.95, 11.6, 1.5, "Evid\u00eancia para uso industrial",
     size=48, font=F_TITLE, bold=True, color=RGBColor(0xF2, 0xF5, 0xF4))
text(s, 0.85, 3.5, 8.6, 0.8,
     "Contra o que o DisperSim 3D j\u00e1 foi testado, o que ele reproduz e onde ele para.",
     size=17, color=RGBColor(0xB9, 0xC7, 0xC7), spacing=1.35)
box(s, 0.85, 4.6, 11.63, 0.012, fill=RGBColor(0x2C, 0x40, 0x42))
text(s, 0.85, 4.95, 11.6, 1.2,
     ["Este doss\u00ea foi escrito para ser conferido. Cada n\u00famero vem de um arquivo de benchmark no",
      "reposit\u00f3rio, e cada arquivo carrega a cita\u00e7\u00e3o do experimento que codifica. As falhas aparecem",
      "com o mesmo destaque dos acertos, porque um modelo de consequ\u00eancias com limites n\u00e3o",
      "documentados n\u00e3o sustenta um caso de seguran\u00e7a."],
     size=12.5, color=RGBColor(0x8E, 0xA3, 0xA3), spacing=1.4)
text(s, 0.85, 6.7, 11.63, 0.3, "Rev. 28/08/2026   \u00b7   branch master",
     size=10, font=F_MONO, color=RGBColor(0x6C, 0x81, 0x82), caps=True, tracking=1.3)

# ════════════════════════════════════════════════════════════════════
# 2. Placar
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Resultado geral", "O placar")
cards = [
    ("Radia\u00e7\u00e3o de inc\u00eandio", "16", "/ 16", GOOD,
     "Po\u00e7as de GNL e de hidrocarbonetos de 6 a 80 m, e jatos de g\u00e1s natural de 2,8 a 19,5 kg/s, todos em escala de campo."),
    ("Autotestes num\u00e9ricos", "271", "/ 271", GOOD,
     "Sete su\u00edtes: geometria, tabela de detectores IOGP, persist\u00eancia, radia\u00e7\u00e3o, flash fire, dose t\u00e9rmica e estudo de inc\u00eandio."),
    ("Dispers\u00e3o de gases", "18", "/ 31", WARN,
     "Contra o desempenho publicado do FLACS e do PHAST nos mesmos ensaios. Seis das treze diferen\u00e7as s\u00e3o fronteiras de escopo declaradas."),
    ("Experimentos codificados", "47", "", INK,
     "31 ensaios de dispers\u00e3o e 16 ensaios de inc\u00eandio, cada um com a fonte prim\u00e1ria registrada no arquivo."),
]
cw, gap = 2.83, 0.09
for i, (lab, big, small, col, note) in enumerate(cards):
    x = 0.85 + i * (cw + gap)
    box(s, x, y, cw, 2.55, fill=SURFACE, line=RULE)
    text(s, x + 0.25, y + 0.28, cw - 0.5, 0.22, lab, size=8.5, font=F_MONO,
         color=FAINT, caps=True, tracking=1.2)
    text(s, x + 0.25, y + 0.62, cw - 0.5, 0.7,
         [[(big, {"size": 44, "color": col}), (" " + small, {"size": 17, "color": MUTED})]],
         font=F_TITLE, bold=True)
    text(s, x + 0.25, y + 1.45, cw - 0.5, 1.0, note, size=10, color=MUTED, spacing=1.28)

text(s, 0.85, y + 2.85, 11.63, 0.9,
     ["Os n\u00fameros de inc\u00eandio e os autotestes foram reexecutados para este doss\u00ea no commit atual de master.",
      "Os n\u00fameros de dispers\u00e3o v\u00eam da rodada registrada pelo pr\u00f3prio projeto em 16/05/2026 e n\u00e3o foram reexecutados aqui."],
     size=11.5, color=MUTED, spacing=1.4)
footer(s, 2)

# ════════════════════════════════════════════════════════════════════
# 3. O que conta como aprovacao
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "M\u00e9todo", "O que conta como aprova\u00e7\u00e3o, e o que n\u00e3o conta")
text(s, 0.85, y, 6.6, 2.6,
     ["Um benchmark compara uma previs\u00e3o com uma medi\u00e7\u00e3o publicada e verifica se a raz\u00e3o entre as duas cai dentro de uma banda declarada. A banda padr\u00e3o \u00e9 a conven\u00e7\u00e3o de fator dois usada em toda a modelagem de consequ\u00eancias: raz\u00e3o prevista/medida entre 0,5 e 2,0.",
      "Para dispers\u00e3o a compara\u00e7\u00e3o vai um passo al\u00e9m. Aprova\u00e7\u00e3o ali n\u00e3o quer dizer que a previs\u00e3o bateu com o experimento. Quer dizer que as medidas estat\u00edsticas de Hanna ficaram dentro de uma toler\u00e2ncia documentada em rela\u00e7\u00e3o \u00e0s que o FLACS ou o PHAST obtiveram no mesmo ensaio."],
     size=13, color=INK, spacing=1.4, space_after=12)

callout(s, 0.85, y + 2.75, 6.6, 1.55, "Falcon 1, o exemplo mais claro",
        ["O DisperSim reporta FAC2 = 0,00 e passa, porque a coorte FLACS publicada para Falcon com barreira de vapor tamb\u00e9m reporta FAC2 = 0,00. Os dois c\u00f3digos sofrem naquele caso. O benchmark estabelece que o DisperSim n\u00e3o \u00e9 pior, e nada al\u00e9m disso."],
        bar=SIGNAL)

x2 = 7.9
text(s, x2, y - 0.05, 4.6, 0.3, "Tr\u00eas classes de evid\u00eancia, mantidas separadas",
     size=13.5, font=F_TITLE, bold=True, color=INK)
classes = [
    ("Ensaio de campo", "Comparado com medi\u00e7\u00f5es de um experimento publicado. \u00c9 a \u00fanica classe capaz de dizer que o modelo acerta o mundo.", ACCENT),
    ("Autoconsist\u00eancia", "O motor conferido contra a pr\u00f3pria solu\u00e7\u00e3o anal\u00edtica. Pega regress\u00e3o num\u00e9rica e n\u00e3o prova nada sobre f\u00edsica.", MUTED),
    ("Linha de base de regress\u00e3o", "O motor conferido contra a pr\u00f3pria \u00faltima sa\u00edda boa conhecida, onde o dado prim\u00e1rio \u00e9 restrito. Nunca apresentada como concord\u00e2ncia com um experimento.", MUTED),
]
yy = y + 0.42
for nome, desc, col in classes:
    box(s, x2, yy, 4.6, 1.28, fill=SURFACE, line=RULE)
    box(s, x2, yy, 0.04, 1.28, fill=col)
    text(s, x2 + 0.28, yy + 0.18, 4.1, 0.25, nome, size=12, font=F_TITLE, bold=True, color=INK)
    text(s, x2 + 0.28, yy + 0.5, 4.1, 0.7, desc, size=10, color=MUTED, spacing=1.3)
    yy += 1.42
footer(s, 3)

# ════════════════════════════════════════════════════════════════════
# 4. dataConfidence
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "M\u00e9todo", "Dado que ningu\u00e9m conferiu nunca \u00e9 aprova\u00e7\u00e3o")
text(s, 0.85, y, 7.4, 1.3,
     ["Todo benchmark de inc\u00eandio declara um dataConfidence. Um benchmark Unverified \u00e9 avaliado e impresso, mas o runner o reporta como n\u00e3o contabilizado em vez de aprovado. Um visto verde contra um n\u00famero que ningu\u00e9m conferiu \u00e9 pior do que n\u00e3o ter teste nenhum."],
     size=14, color=INK, spacing=1.42)

table(s, 0.85, y + 1.45, 7.4,
      ["N\u00edvel", "Significado", "Nos 16 ensaios"],
      [[("High", {"font": F_MONO, "bold": True, "color": GOOD}),
        "Lido de uma tabela na fonte citada",
        ("7", {"font": F_MONO, "align": PP_ALIGN.RIGHT})],
       [("Medium", {"font": F_MONO, "bold": True, "color": WARN}),
        "Lido de uma figura, ou de uma cita\u00e7\u00e3o secund\u00e1ria",
        ("9", {"font": F_MONO, "align": PP_ALIGN.RIGHT})],
       [("Unverified", {"font": F_MONO, "bold": True, "color": BAD}),
        "N\u00e3o conferido contra a fonte. Nunca contabilizado",
        ("0", {"font": F_MONO, "align": PP_ALIGN.RIGHT})]],
      widths=[1.5, 5.0, 1.6], rowh=0.42, fsize=11.5,
      aligns=[PP_ALIGN.LEFT, PP_ALIGN.LEFT, PP_ALIGN.RIGHT])

callout(s, 8.55, y, 3.93, 3.25, "Como a su\u00edte roda",
        ["Cada arquivo .fbench descreve o ensaio em duas metades e o runner pontua as duas separado.",
         "1. Geometria da chama e poder emissivo, que testam as correla\u00e7\u00f5es e o balan\u00e7o de energia.",
         "2. Fluxo incidente nos radi\u00f4metros, que testa o fator de forma e a transmissividade em cima delas.",
         "Saber qual metade falhou \u00e9 a maior parte do conserto."])
footer(s, 4)

# ════════════════════════════════════════════════════════════════════
# 5. Pocas de GNL
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Ensaio de campo \u00b7 Radia\u00e7\u00e3o", "Po\u00e7as de GNL: poder emissivo superficial",
           "Raj (2005) Tabela 1  \u00b7  5 ensaios  \u00b7  confian\u00e7a m\u00e9dia  \u00b7  di\u00e2metros de 6 a 36 m")
gauge_block(s, 0.85, y, 11.6, [
    ("AGA San Clemente", "6,1 m  \u00b7  1973", "143 a 178 kW/m\u00b2", 0.263, "0,72", WARN),
    ("China Lake, sobre \u00e1gua", "15 m  \u00b7  1974 a 1976", "220 \u00b1 30 kW/m\u00b2", 0.408, "0,88", GOOD),
    ("Esso L\u00edbia, trincheira", "18 m  \u00b7  1969", "92 kW/m\u00b2", 0.391, "0,86", GOOD),
    ("Maplin Sands", "20 m  \u00b7  1980", "150 a 220 kW/m\u00b2", 0.391, "0,86", GOOD),
    ("Montoir", "35,7 m  \u00b7  1987", "257 a 273 kW/m\u00b2", 0.400, "0,87", GOOD),
    ("Montoir, comprimento", "o \u00fanico publicado aqui", "78 m", 0.400, "0,87", GOOD),
], label_w=3.3, obs_w=2.0, track_w=4.6)

box(s, 0.85, 6.25, 11.63, 0.012, fill=RULE)
text(s, 0.85, 6.42, 11.6, 0.4,
     "Di\u00e2metros de po\u00e7a variando seis vezes, poderes emissivos variando tr\u00eas vezes. Todos dentro da banda.",
     size=12, color=MUTED)
footer(s, 5)

# ════════════════════════════════════════════════════════════════════
# 6. O vies de 13%
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Ensaio de campo \u00b7 Radia\u00e7\u00e3o", "O vi\u00e9s \u00e9 o achado, n\u00e3o os cinco vistos verdes")

box(s, 0.85, y, 5.3, 2.05, fill=SURFACE, line=RULE)
text(s, 1.2, y + 0.3, 4.6, 0.3, "Quatro das cinco raz\u00f5es", size=10, font=F_MONO,
     color=FAINT, caps=True, tracking=1.3)
text(s, 1.2, y + 0.65, 4.6, 0.9, "0,86 a 0,88", size=48, font=F_TITLE, bold=True, color=SIGNAL)
text(s, 1.2, y + 1.5, 4.6, 0.4,
     "Subestima\u00e7\u00e3o sistem\u00e1tica de 13% ao longo de uma varia\u00e7\u00e3o de seis vezes em di\u00e2metro.",
     size=11, color=MUTED, spacing=1.3)

text(s, 6.55, y + 0.02, 5.93, 2.1,
     ["O modelo acompanha como o poder emissivo escala com o di\u00e2metro da po\u00e7a e ent\u00e3o subestima o n\u00edvel de forma consistente. Isso n\u00e3o \u00e9 dispers\u00e3o estat\u00edstica.",
      "O lugar mais prov\u00e1vel para o vi\u00e9s \u00e9 a fra\u00e7\u00e3o radiativa \u03c7, que n\u00e3o \u00e9 reportada nem para Montoir nem para Maplin Sands. 0,25 foi assumido para as duas, e um valor perto de 0,29 colocaria ambas em cima do alvo."],
     size=13, color=INK, spacing=1.4, space_after=11)

callout(s, 0.85, y + 2.35, 11.63, 1.95, "N\u00e3o foi ajustado, de prop\u00f3sito",
        ["Ajustar uma constante a cinco pontos n\u00e3o \u00e9 valida\u00e7\u00e3o, e separar \u03c7 da correla\u00e7\u00e3o de comprimento de chama exige dados de radi\u00f4metro que esses cinco ensaios n\u00e3o t\u00eam.",
         "Os 13% ficam declarados para quem revisa decidir se importam no caso dele, lembrando que o desvio \u00e9 conservador para comprimento de chama e n\u00e3o conservador para fluxo."],
        bar=SIGNAL)
footer(s, 6)

# ════════════════════════════════════════════════════════════════════
# 7. Pocas fuliginosas
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Ensaio de campo \u00b7 Radia\u00e7\u00e3o", "Po\u00e7as de hidrocarbonetos fuliginosos",
           "Mudan (1984) Tabela 4  \u00b7  4 ensaios  \u00b7  confian\u00e7a m\u00e9dia")
y = y - 0.3
yy = gauge_block(s, 0.85, y, 11.6, [
    ("Gasolina (Fu)", "10 m", "60 a 130 kW/m\u00b2 m\u00e1x", 0.455, "0,94", GOOD),
    ("GLP em terra (Mizner e Eyre)", "20 m", "48 kW/m\u00b2", 0.178, "0,64", WARN),
    ("Querosene (JISE)", "30 m", "10 a 25 kW/m\u00b2 m\u00e9d", 0.706, "1,33", GOOD),
    ("Querosene (JISE)", "80 m", "10 a 25 kW/m\u00b2 m\u00e9d", 0.595, "1,14", GOOD),
], label_w=3.3, obs_w=2.0, track_w=4.6)

text(s, 0.85, yy + 0.2, 5.6, 1.2,
     "Os dois casos de querosene s\u00e3o o motivo do exerc\u00edcio: a mistura de fuligem achata com o di\u00e2metro do mesmo jeito que as medi\u00e7\u00f5es, 1,33 aos 30 m caindo para 1,14 aos 80 m, que \u00e9 a obscura\u00e7\u00e3o por fuma\u00e7a que ela existe para representar.",
     size=11.5, color=INK, spacing=1.35)

callout(s, 6.9, yy + 0.12, 5.58, 1.55, "O erro de 0,64 no GLP ficou como est\u00e1",
        ["A Tabela 2 de Mudan d\u00e1 uma fra\u00e7\u00e3o radiativa medida de 7% para esse ensaio, e o balan\u00e7o de energia puro com ela chega a 47 kW/m\u00b2 contra 48 medidos. \u00c9 o teto de fuligem que puxa a previs\u00e3o para 31."],
        bar=WARN, size=11)
footer(s, 7)

# ════════════════════════════════════════════════════════════════════
# 8. Jatos: comprimento de chama
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Ensaio de campo \u00b7 Radia\u00e7\u00e3o", "Jatos de g\u00e1s natural: comprimento de chama",
           "Miller (2017), casos DNV GL  \u00b7  3 ensaios  \u00b7  confian\u00e7a alta")
yy = gauge_block(s, 0.85, y, 11.6, [
    ("2,9 kg/s horizontal", "", "19,8 m", 0.613, "1,17", GOOD),
    ("9,6 kg/s horizontal", "", "37,8 m", 0.493, "0,99", GOOD),
    ("19,5 kg/s horizontal", "", "49,9 m", 0.500, "1,00", GOOD),
], label_w=3.3, obs_w=2.0, track_w=4.6)

box(s, 0.85, yy + 0.35, 11.63, 0.012, fill=RULE)
text(s, 0.85, yy + 0.6, 5.6, 0.9,
     "A faixa de vazamento que este projeto atende \u00e9 de 0,1 a 32 kg/s por orif\u00edcios de 1 a 50 mm. Estes tr\u00eas casos ficam dentro dela e cobrem a metade superior.",
     size=13, color=INK, spacing=1.4)

callout(s, 6.9, yy + 0.5, 5.58, 1.7, "Um defeito achado por autoteste",
        ["A correla\u00e7\u00e3o de Chamberlain L = 0,2 Q^0,4 recebe a taxa de libera\u00e7\u00e3o de calor em quilowatts e estava recebendo watts. Um jato de 2,9 kg/s sa\u00eda com centenas de metros, 16 vezes o comprimento real."],
        bar=BAD)
footer(s, 8)

# ════════════════════════════════════════════════════════════════════
# 9. Jatos: fluxo nos radiometros
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Ensaio de campo \u00b7 Radia\u00e7\u00e3o", "Jatos de g\u00e1s natural: fluxo nos radi\u00f4metros",
           "Johnson (1994), tabulado por Miller (2017)  \u00b7  4 ensaios, 30 radi\u00f4metros  \u00b7  confian\u00e7a alta")
gauge_block(s, 0.85, y, 11.6, [
    ("Ensaio 1033", "7,9 kg/s \u00b7 75 mm \u00b7 11,1 barg \u00b7 vento a 1\u00b0", "8 radi\u00f4metros", (0.562, 0.695), "1,09 a 1,31", GOOD),
    ("Ensaio 1040", "2,8 kg/s \u00b7 152 mm \u00b7 0,3 barg \u00b7 vento a \u221223\u00b0", "7 radi\u00f4metros", (0.471, 0.866), "0,96 a 1,66", GOOD),
    ("Ensaio 1089", "3,8 kg/s \u00b7 20 mm \u00b7 66 barg \u00b7 vento a \u22121\u00b0", "5 radi\u00f4metros", (0.738, 0.773), "1,39 a 1,46", GOOD),
    ("Ensaio 1083, a jusante", "8,4 kg/s \u00b7 152 mm \u00b7 2,1 barg \u00b7 vento a 56\u00b0", "10 radi\u00f4metros", (0.167, 0.243), "0,63 a 0,70", GOOD),
    ("Ensaio 1083, R10", "44 m ao lado da chama", "transversal", 0.98, "3,67 \u2192", BAD),
], label_w=3.3, obs_w=2.0, track_w=4.6)

box(s, 0.85, 5.95, 11.63, 0.012, fill=RULE)
text(s, 0.85, 6.15, 11.6, 0.6,
     [[("O ensaio 1033 \u00e9 aquele cuja geometria o modelo de fato representa", {"bold": True}),
       (", com o vento alinhado ao vazamento, e ele fica dentro de 31% em todos os oito radi\u00f4metros. A\u00ed est\u00e3o o fator de forma, a transmissividade e o poder emissivo concordando com a medi\u00e7\u00e3o ao mesmo tempo.", {})]],
     size=12, color=INK, spacing=1.35)
footer(s, 9)

# ════════════════════════════════════════════════════════════════════
# 10. Onde os outros tres degradam
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Limita\u00e7\u00e3o nomeada", "Onde os outros tr\u00eas degradam, e por qu\u00ea")
text(s, 0.85, y, 11.6, 0.4,
     "Os outros tr\u00eas degradam de um jeito que diz o que est\u00e1 faltando, que \u00e9 a parte \u00fatil.",
     size=13.5, color=INK)

box(s, 0.85, y + 0.55, 5.7, 2.35, fill=SURFACE, line=RULE)
box(s, 0.85, y + 0.55, 0.04, 2.35, fill=WARN)
text(s, 1.15, y + 0.78, 5.15, 0.3, "1089 superestima de 39 a 46%", size=14, font=F_TITLE, bold=True, color=INK)
text(s, 1.15, y + 1.18, 5.15, 1.6,
     "\u00c9 o caso de alta press\u00e3o, 66 barg por um furo de 20 mm. Miller registra que o modelo AP Flame superestima esse mesmo ensaio, ent\u00e3o \u00e9 um caso reconhecidamente dif\u00edcil e n\u00e3o uma peculiaridade desta implementa\u00e7\u00e3o. Superestimar fluxo \u00e9 o lado conservador.",
     size=11.5, color=MUTED, spacing=1.35)

box(s, 6.78, y + 0.55, 5.7, 2.35, fill=SURFACE, line=RULE)
box(s, 6.78, y + 0.55, 0.04, 2.35, fill=BAD)
text(s, 7.08, y + 0.78, 5.15, 0.3, "1040 e 1083 t\u00eam o vento fora do eixo", size=14, font=F_TITLE, bold=True, color=INK)
text(s, 7.08, y + 1.18, 5.15, 1.6,
     "A 23\u00b0 e 56\u00b0 do eixo do vazamento. O modelo inclina uma chama de eixo \u00fanico na dire\u00e7\u00e3o do vento e a levanta por empuxo no plano vertical que cont\u00e9m esse eixo, mas n\u00e3o consegue curv\u00e1-la lateralmente. Os radi\u00f4metros ao longo da chama continuam precisos; os transversais superestimam.",
     size=11.5, color=MUTED, spacing=1.35)

callout(s, 0.85, y + 3.1, 11.63, 1.3, "O R10 do ensaio 1083 fica reportado a 3,67\u00d7",
        ["A 44 m do lado da chama, ele \u00e9 a posi\u00e7\u00e3o mais exposta exatamente a essa limita\u00e7\u00e3o. Alargar a banda de aceita\u00e7\u00e3o para deixar a su\u00edte verde jogaria fora uma afirma\u00e7\u00e3o verdadeira sobre o modelo."],
        bar=BAD)
footer(s, 10)

# ════════════════════════════════════════════════════════════════════
# 11. Arqueamento por empuxo
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Ensaio de campo \u00b7 Radia\u00e7\u00e3o", "Arqueamento por empuxo de jatos horizontais",
           "Miller (2017), equa\u00e7\u00f5es (19) a (22)  \u00b7  n\u00famero de Richardson na fonte expandida")
text(s, 0.85, y, 11.6, 0.85,
     "Uma chama-jato horizontal n\u00e3o fica no eixo do vazamento. Ela segue reta enquanto o pr\u00f3prio momento domina, e depois o empuxo vira o restante para cima. A divis\u00e3o entre as duas se\u00e7\u00f5es \u00e9 dada pelo n\u00famero de Richardson avaliado na fonte expandida e n\u00e3o no furo: um orif\u00edcio de 20 mm a 66 barg e um de 152 mm a 0,3 barg podem passar a mesma massa e produzir chamas completamente diferentes.",
     size=13, color=INK, spacing=1.4)

table(s, 0.85, y + 1.0, 11.63,
      ["Ensaio", "A montante", "Raz\u00f5es de fluxo antes", "Raz\u00f5es de fluxo depois", "Saldo"],
      [[("1033", {"font": F_MONO}), ("11,1 barg", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        ("0,96 a 1,18", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        ("1,09 a 1,31", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        "um pouco pior, a chama j\u00e1 era quase reta"],
       [("1040", {"font": F_MONO}), ("0,3 barg", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        ("0,99 a 1,86", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        ("0,96 a 1,66", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        ("melhor", {"color": GOOD})],
       [("1083", {"font": F_MONO}), ("2,1 barg", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        ("0,43 a 0,45", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        ("0,63 a 0,70", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        ("falha \u2192 passa", {"color": GOOD, "bold": True})],
       [("1089", {"font": F_MONO}), ("66 barg", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        ("1,28 a 1,43", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        ("1,39 a 1,46", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        "um pouco pior, a chama j\u00e1 era quase reta"]],
      widths=[1.0, 1.3, 1.9, 1.9, 4.3], rowh=0.42, fsize=11.5)

text(s, 0.85, y + 3.05, 11.6, 0.55,
     "Um FireSource sem press\u00e3o de estagna\u00e7\u00e3o \u00e9 tratado como subs\u00f4nico e a chama dele fica reta, ent\u00e3o todo projeto salvo antes dessa mudan\u00e7a mant\u00e9m o comportamento anterior.",
     size=11.5, color=MUTED, spacing=1.35)
footer(s, 11)

# ════════════════════════════════════════════════════════════════════
# 12. Quatro defeitos
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Por que isso importa", "Quatro defeitos que a literatura achou",
           "Uma taxa de aprova\u00e7\u00e3o alta diz pouco sozinha. O que sustenta um modelo de consequ\u00eancias em uso industrial \u00e9 a prova de que a compara\u00e7\u00e3o com medi\u00e7\u00f5es reais consegue reprov\u00e1-lo, e nenhum destes quatro seria pego por um teste unit\u00e1rio.")
y = y - 0.22

defects = [
    ("Montoir 35 m GNL", "Teto de fuligem em combust\u00edvel limpo",
     "O teto de Mudan estava sendo aplicado a toda po\u00e7a. Ele \u00e9 calibrado em hidrocarbonetos fuliginosos e limita o poder emissivo perto de 22 kW/m\u00b2. O GNL queima limpo: os ensaios reportam de 165 a 265. O primeiro benchmark previu 21,8 contra 265 medidos."),
    ("Johnson 1083", "Quem estava na frente recebia zero",
     "A chama era painelada s\u00f3 na superf\u00edcie lateral. Um receptor em frente \u00e0 ponta v\u00ea todo painel de perfil ou por tr\u00e1s, e a soma d\u00e1 exatamente zero. Tr\u00eas radi\u00f4metros mediram 4,6, 3,3 e 2,2 kW/m\u00b2 e o modelo devolvia 0."),
    ("A coorte de GNL", "Normaliza\u00e7\u00e3o pela \u00e1rea errada",
     "Incluir a tampa da ponta na \u00e1rea emissora derrubou o SEP de toda po\u00e7a em 13% sem f\u00edsica por tr\u00e1s. Os valores publicados s\u00e3o normalizados pela superf\u00edcie lateral do cilindro, e o modelo agora cita a mesma grandeza que a medi\u00e7\u00e3o cita."),
    ("Autoteste de campo distante", "Chamas 16 vezes maiores",
     "A correla\u00e7\u00e3o de Chamberlain recebe a taxa de libera\u00e7\u00e3o de calor em quilowatts e estava recebendo watts. Um jato de 2,9 kg/s sa\u00eda com centenas de metros. Esse foi pego por um autoteste, que \u00e9 no que autoteste \u00e9 bom."),
]
cw = 2.83
for i, (achado, titulo, desc) in enumerate(defects):
    x = 0.85 + i * (cw + 0.09)
    box(s, x, y + 0.72, cw, 3.55, fill=SURFACE, line=RULE)
    text(s, x + 0.22, y + 0.95, cw - 0.44, 0.5, "Achado por\n" + achado, size=8.5,
         font=F_MONO, color=BAD, caps=True, tracking=1.0, spacing=1.3)
    text(s, x + 0.22, y + 1.55, cw - 0.44, 0.7, titulo, size=13.5, font=F_TITLE,
         bold=True, color=INK, spacing=1.1)
    text(s, x + 0.22, y + 2.35, cw - 0.44, 1.8, desc, size=10, color=MUTED, spacing=1.32)
footer(s, 12)

# ════════════════════════════════════════════════════════════════════
# 13. Autotestes
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Autoconsist\u00eancia", "Autotestes num\u00e9ricos: 271 asser\u00e7\u00f5es",
           "Estes n\u00e3o estabelecem que a f\u00edsica est\u00e1 certa. Estabelecem que a aritm\u00e9tica, a geometria e a persist\u00eancia n\u00e3o derivam em sil\u00eancio.")
y = y - 0.3
table(s, 0.85, y, 11.63,
      ["Su\u00edte", "Asser\u00e7\u00f5es", "O que ela trava", "Flag"],
      [["Geometria", ("19", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        "Voxeliza\u00e7\u00e3o, limites de obst\u00e1culo, transforma\u00e7\u00f5es de coordenadas",
        ("--geometry-selftest", {"font": F_MONO, "size": 9.5, "color": MUTED})],
       ["Tabela de detectores IOGP", ("27", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        "Tabela IOGP 434-01 embutida contra os valores publicados",
        ("--iogp-selftest", {"font": F_MONO, "size": 9.5, "color": MUTED})],
       ["Ida e volta do cen\u00e1rio de inc\u00eandio", ("111", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        "Todo campo de fogo, igni\u00e7\u00e3o e estudo sobrevive a salvar e recarregar",
        ("--fire-roundtrip-selftest", {"font": F_MONO, "size": 9.5, "color": MUTED})],
       ["Radia\u00e7\u00e3o de chama s\u00f3lida", ("36", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        "Fator de forma contra o limite anal\u00edtico de fonte pontual, transmissividade, SEP",
        ("--solid-flame-selftest", {"font": F_MONO, "size": 9.5, "color": MUTED})],
       ["Flash fire", ("17", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        "Envelope LFL a UFL, componentes conexas, tempos de chegada da queima geod\u00e9sica",
        ("--flash-fire-selftest", {"font": F_MONO, "size": 9.5, "color": MUTED})],
       ["Dose t\u00e9rmica", ("29", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        "Integral de dose, probit de fatalidade de Eisenberg, probits de queimadura, a erf",
        ("--thermal-dose-selftest", {"font": F_MONO, "size": 9.5, "color": MUTED})],
       ["Estudo de inc\u00eandio", ("32", {"font": F_MONO, "align": PP_ALIGN.RIGHT}),
        "Pontua\u00e7\u00e3o de fontes, ranking de igni\u00e7\u00f5es, montagem da cena de fonte \u00fanica",
        ("--fire-study-selftest", {"font": F_MONO, "size": 9.5, "color": MUTED})]],
      widths=[2.9, 1.1, 5.1, 2.6], rowh=0.36, fsize=10.5)

callout(s, 0.85, y + 3.15, 11.63, 1.25, "Um exemplo do que eles pegam",
        ["A su\u00edte de chama s\u00f3lida afirma que um cilindro painelado visto de longe converge para o resultado de fonte pontual. A primeira vers\u00e3o dessa afirma\u00e7\u00e3o esperava raz\u00e3o 1,0, quando um cilindro visto de lado irradia 4/\u03c0 \u2248 1,273 vezes o fluxo isotr\u00f3pico. O teste agora afirma a anisotropia em vez de encobri-la."])
footer(s, 13)

# ════════════════════════════════════════════════════════════════════
# 14. Dispersao por coorte
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Ensaio de campo \u00b7 Dispers\u00e3o", "Dispers\u00e3o: 18 aprovados em 31",
           "Rodada registrada em 16/05/2026  \u00b7  o n\u00famero s\u00f3 \u00e9 leg\u00edvel depois de separar as treze diferen\u00e7as por causa")
y = y - 0.3
table(s, 0.85, y, 11.63,
      ["Coorte", "Passa", "Total", "Refer\u00eancia contra a qual a aprova\u00e7\u00e3o \u00e9 medida"],
      [["Autoconsist\u00eancia (pluma e sopro gaussianos)", ("2", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": GOOD, "bold": True}), ("2", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "Solu\u00e7\u00e3o anal\u00edtica do pr\u00f3prio motor"],
       ["Prairie Grass, terreno plano, neutro a est\u00e1vel", ("3", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": GOOD, "bold": True}), ("5", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "Coorte FLACS, Hanna e Chang 2004"],
       ["Campo de GNL: Burro, Coyote, Maplin Sands", ("5", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": GOOD, "bold": True}), ("8", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "Coorte FLACS sem obstru\u00e7\u00e3o, Hansen 2010"],
       ["Falcon, GNL com barreira de vapor", ("3", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": GOOD, "bold": True}), ("3", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "Coorte FLACS Falcon, Hansen 2010"],
       ["SF\u2086 em t\u00fanel de vento sobre rampa (DAT632)", ("1", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": GOOD, "bold": True}), ("1", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "Linha de base de regress\u00e3o"],
       ["MUST, arranjo urbano de 120 obst\u00e1culos", ("1", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": GOOD, "bold": True}), ("1", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "Linha de base de regress\u00e3o (dado por ensaio \u00e9 restrito)"],
       ["Jatos no FluidX3D: CH\u2084, CO\u2082, H\u2082", ("3", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": GOOD, "bold": True}), ("5", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "Linhas de base de regress\u00e3o e Witlox 2014"],
       [("G\u00e1s denso no motor gaussiano", {"color": MUTED}), ("0", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": WARN, "bold": True}), ("6", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": MUTED}), ("Fora de escopo por declara\u00e7\u00e3o", {"color": MUTED})]],
      widths=[4.3, 0.9, 0.9, 5.5], rowh=0.35, fsize=10.5)

callout(s, 0.85, y + 3.22, 11.63, 1.15, "Onde o DisperSim \u00e9 melhor que a refer\u00eancia",
        ["O desempenho publicado do FLACS em Falcon com barreira de vapor \u00e9 FAC2 = 0,00, MG = 5,56, VG = 23,65. O DisperSim chega a FAC2 = 1,00 em Falcon 3 e Falcon 4, com MG de 1,40 e 1,26."],
        bar=GOOD)
footer(s, 14)

# ════════════════════════════════════════════════════════════════════
# 15. As 13 diferencas
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Fronteira do envelope", "As treze diferen\u00e7as, separadas por causa")

text(s, 0.85, y, 5.7, 0.55,
     [[("Seis fronteiras de escopo declaradas. ", {"bold": True, "color": WARN}),
       ("Rodam o motor gaussiano contra f\u00edsica de g\u00e1s denso que ele nunca foi feito para conter, e ficam na su\u00edte de prop\u00f3sito.", {})]],
     size=11.5, color=INK, spacing=1.3)
table(s, 0.85, y + 0.72, 5.7,
      ["Ensaio", "F\u00edsica ausente"],
      [["Kit Fox U5-2", "G\u00e1s denso em arranjo de obst\u00e1culos"],
       ["Desert Tortoise 4", "Aerossol e rainout"],
       ["Thorney Island 8", "Espalhamento gravitacional"],
       ["Jack Rabbit I, 7", "Sa\u00edda de depress\u00e3o de terreno"],
       ["Jack Rabbit II, 1", "Arranjo urbano"],
       ["Jack Rabbit II, 7", "Campo distante, erra log MG por 0,03"]],
      widths=[2.1, 3.6], rowh=0.36, fsize=10.5)

text(s, 6.78, y, 5.7, 0.55,
     [[("Sete limita\u00e7\u00f5es reais em aberto. ", {"bold": True, "color": BAD}),
       ("Duas s\u00e3o decis\u00f5es apertadas contra toler\u00e2ncia apertada; tr\u00eas s\u00e3o o Sc\u209c de f\u00e1brica; duas s\u00e3o f\u00edsica ausente.", {})]],
     size=11.5, color=INK, spacing=1.3)
table(s, 6.78, y + 0.72, 5.7,
      ["Ensaio", "MRB", "FAC2", "Erro"],
      [["Prairie Grass 7 (B)", ("\u22120,69", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("0,60", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "|MRB| por 0,01"],
       ["Prairie Grass 11 (C)", ("0,72", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("0,40", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "|MRB| por 0,04"],
       ["Burro 9", ("0,72", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("0,67", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "MG por 0,23, Sc\u209c"],
       ["Coyote 3", ("0,76", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("0,60", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "FAC2 por 0,04, Sc\u209c"],
       ["Maplin Sands 27", ("1,41", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("0,00", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), "Sc\u209c"],
       ["CO2PipeHaz 6 mm", ("n/d", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": FAINT}), ("n/d", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": FAINT}), "Sem modelo bif\u00e1sico"],
       ["Jato de H\u2082 (Schefer)", ("n/d", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": FAINT}), ("n/d", {"font": F_MONO, "align": PP_ALIGN.RIGHT, "color": FAINT}), "Sem dado por sensor"]],
      widths=[2.1, 0.85, 0.85, 1.9], rowh=0.36, fsize=10.5)

text(s, 0.85, y + 3.35, 11.63, 0.75,
     "O caso do hidrog\u00eanio \u00e9 contabilidade honesta: o motor roda limpo e produz uma nuvem plaus\u00edvel, mas n\u00e3o foi encontrado dado publicado por sensor para aquele vazamento, ent\u00e3o a compara\u00e7\u00e3o no arquivo n\u00e3o tem autoridade e ela \u00e9 reportada como falha em vez de ser afrouxada em sil\u00eancio.",
     size=11.5, color=MUTED, spacing=1.35)
footer(s, 15)

# ════════════════════════════════════════════════════════════════════
# 16. Resultado negativo de Vu
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Resultado negativo", "O Schmidt turbulento, e um resultado registrado",
           "As tr\u00eas falhas de GNL no OpenFOAM s\u00e3o casos em que Vu (2019) chega a FAC2 = 1,00 com um solver pr\u00f3prio. A receita foi implementada e medida no Coyote 3.")
y = y - 0.3
table(s, 0.85, y, 11.63,
      ["Pilha", "MRB", "FAC2", "MG", "Veredito"],
      [["rhoReactingBuoyantFoam de f\u00e1brica", ("0,76", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("0,60", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("2,30", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("linha de base", {"color": MUTED})],
       ["+ solver corrigido com Sc\u209c = 0,15", ("0,72", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("0,40", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("2,15", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("pior", {"color": WARN})],
       ["+ refino de malha, 3 caixas alinhadas ao vento", ("1,44", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("0,00", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("6,37", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("muito pior", {"color": BAD})],
       ["+ precursor de CLA, 500 itera\u00e7\u00f5es SIMPLE", ("1,78", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("0,00", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("17,75", {"font": F_MONO, "align": PP_ALIGN.RIGHT}), ("muito pior de novo", {"color": BAD, "bold": True})]],
      widths=[5.0, 1.2, 1.2, 1.2, 3.0], rowh=0.42, fsize=11)

callout(s, 0.85, y + 2.28, 11.63, 1.75, "Toda modifica\u00e7\u00e3o degradou a previs\u00e3o de forma monot\u00f4nica",
        ["A linha de base daqui j\u00e1 subestima onde a de Vu superestima, ent\u00e3o as mudan\u00e7as dela, que amplificam difus\u00e3o, empurram no sentido errado para este pipeline.",
         "Os tr\u00eas flags existem na configura\u00e7\u00e3o e os tr\u00eas est\u00e3o desligados por padr\u00e3o, com esta tabela como motivo. A diferen\u00e7a que sobra \u00e9 atribu\u00edda a algo a montante do solver, provavelmente a condi\u00e7\u00e3o de contorno da fonte, e essa auditoria \u00e9 trabalho em aberto."],
        bar=BAD)
footer(s, 16)

# ════════════════════════════════════════════════════════════════════
# 17. Limitacoes declaradas
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Leia antes de usar", "Limita\u00e7\u00f5es declaradas",
           "Um modelo de consequ\u00eancias serve para uso industrial quando se sabe onde ele vale. Este \u00e9 o envelope atual, escrito sem rodeio.")
y = y - 0.28
lims = [
    ("Curvatura lateral da chama pelo vento n\u00e3o \u00e9 modelada",
     "Receptores ao lado da chama s\u00e3o superestimados, com 3,67\u00d7 na pior posi\u00e7\u00e3o do Johnson 1083. Conservador para loca\u00e7\u00e3o, mas n\u00e3o preciso."),
    ("Subestima\u00e7\u00e3o sistem\u00e1tica de 13% no poder emissivo de po\u00e7as de GNL",
     "Consistente ao longo de uma faixa de seis vezes em di\u00e2metro. N\u00e3o ajustada, por escolha."),
    ("N\u00e3o h\u00e1 modelo bif\u00e1sico, de aerossol nem de rainout na dispers\u00e3o",
     "Am\u00f4nia, CO\u2082 supercr\u00edtico e outros vazamentos com flash est\u00e3o fora do envelope validado."),
    ("O motor gaussiano n\u00e3o serve para g\u00e1s denso, obst\u00e1culos ou depress\u00f5es",
     "Seis benchmarks documentam isso e falham de prop\u00f3sito."),
    ("A mistura de fuligem n\u00e3o est\u00e1 calibrada para GLP",
     "A 20 m ela subestima em 36%, e o balan\u00e7o de energia sem o teto \u00e9 a resposta melhor para esse combust\u00edvel."),
    ("O rhoReactingBuoyantFoam de f\u00e1brica fixa Sc\u209c = 1,0 no c\u00f3digo",
     "Afeta tr\u00eas benchmarks de GNL. Existe um bin\u00e1rio corrigido, desligado por padr\u00e3o."),
    ("O tracer em GPU roda em precis\u00e3o simples",
     "Em jato s\u00f4nico o erro de FP32 chega de 15 a 25% na linha de centro. Use a CPU onde a precis\u00e3o de campo pr\u00f3ximo governa."),
]
yy = y
for i, (t, d) in enumerate(lims):
    col = 0 if i < 4 else 1
    row = i if i < 4 else i - 4
    x = 0.85 + col * 5.93
    ty = y + row * 1.08
    box(s, x, ty, 5.7, 1.0, fill=SURFACE, line=RULE)
    box(s, x, ty, 0.035, 1.0, fill=BAD if i < 2 else WARN)
    text(s, x + 0.26, ty + 0.13, 5.25, 0.34, t, size=11, font=F_TITLE, bold=True,
         color=INK, spacing=1.12)
    text(s, x + 0.26, ty + 0.5, 5.25, 0.44, d, size=9, color=MUTED, spacing=1.26)
footer(s, 17)

# ════════════════════════════════════════════════════════════════════
# 18. Reprodutibilidade
# ════════════════════════════════════════════════════════════════════
s = slide()
y = header(s, "Reprodutibilidade", "Como reproduzir cada n\u00famero deste deck")
text(s, 0.85, y, 6.9, 0.85,
     "As defini\u00e7\u00f5es de benchmark s\u00e3o texto puro no reposit\u00f3rio e o runner faz parte da CLI distribu\u00edda. Nada aqui depende de uma build da m\u00e1quina dos autores nem de um conjunto de dados privado. C\u00f3digo de sa\u00edda 0 significa que toda m\u00e9trica caiu dentro da banda de aceita\u00e7\u00e3o declarada.",
     size=13, color=INK, spacing=1.4)

box(s, 0.85, y + 1.0, 6.9, 2.55, fill=SURFACE, line=RULE)
text(s, 1.1, y + 1.22, 6.4, 2.2,
     ["DisperSim3D.CLI --validate-fire benchmarks/fire",
      "DisperSim3D.CLI --validate      benchmarks/",
      " ",
      "DisperSim3D.CLI --geometry-selftest          # 19",
      "DisperSim3D.CLI --iogp-selftest              # 27",
      "DisperSim3D.CLI --fire-roundtrip-selftest    # 111",
      "DisperSim3D.CLI --solid-flame-selftest       # 36",
      "DisperSim3D.CLI --flash-fire-selftest        # 17",
      "DisperSim3D.CLI --thermal-dose-selftest      # 29",
      "DisperSim3D.CLI --fire-study-selftest        # 32"],
     size=10.5, font=F_MONO, color=INK, spacing=1.35)

table(s, 8.1, y, 4.38,
      ["Ambiente da rodada registrada", ""],
      [["Data", ("16/05/2026", {"font": F_MONO, "size": 9.5})],
       ["CPU", ("i9-13900KF, 24C / 32T", {"font": F_MONO, "size": 9.5})],
       ["GPU", ("RTX 5070 + RTX 3060", {"font": F_MONO, "size": 9.5})],
       ["Runtime", (".NET 10.0.300", {"font": F_MONO, "size": 9.5})],
       ["OpenFOAM", ("v2512 nativo, v2412 WSL2", {"font": F_MONO, "size": 9.5})],
       ["FluidX3D", ("build CUDA nativo", {"font": F_MONO, "size": 9.5})]],
      widths=[1.5, 2.88], rowh=0.36, fsize=10.5)

text(s, 8.1, y + 2.6, 4.38, 1.0,
     "Cada arquivo .dsbench e .fbench guarda as condi\u00e7\u00f5es, a meteorologia, as posi\u00e7\u00f5es dos sensores, os valores observados, a banda de aceita\u00e7\u00e3o e a cita\u00e7\u00e3o da fonte prim\u00e1ria com a tabela ou figura de onde os n\u00fameros vieram.",
     size=10.5, color=MUTED, spacing=1.32)

text(s, 0.85, y + 3.75, 11.63, 0.4,
     "O empacotamento \u00e9 constru\u00eddo e assinado continuamente para Windows, Debian e macOS a partir do mesmo commit que roda estes benchmarks.",
     size=11.5, color=MUTED)
footer(s, 18)

# ════════════════════════════════════════════════════════════════════
# 19-20. Referencias
# ════════════════════════════════════════════════════════════════════
refs_fogo = [
    ("1", "Johnson, A.D., Brightwell, H.M., Carsley, A.J. (1994). Thermal radiation hazards from large scale horizontal natural gas jet fires. Hazards XII, IChemE.", "Ensaios 1033, 1040, 1083, 1089"),
    ("2", "Nedelka, D., Moorhouse, J., Tucker, R.F. (1989). The Montoir 35 m diameter LNG pool fire experiments.", "Montoir 35,7 m"),
    ("3", "Raj, P.K., Atallah, S. (1974). AGA San Clemente LNG fire tests.", "AGA San Clemente 6,1 m"),
    ("4", "Raj, P.K. et al. (1979). USCG China Lake LNG spill and fire tests.", "China Lake 15 m"),
    ("5", "Mizner, G.A., Eyre, J.A. (1983). Large-scale LNG and LPG pool fires.", "Maplin Sands 20 m; GLP 20 m"),
    ("6", "May, W.G., McQueen, W. (1973). Esso LNG trench fire tests, Libya.", "Esso L\u00edbia 18 m"),
    ("7", "Miller, D. (2017). New model for predicting thermal radiation from flares and high pressure jet fires. Process Safety Progress 36(3). DOI 10.1002/prs.11867.", "Fonte expandida, forma da chama, tabelas Johnson, DNV GL"),
    ("8", "Mudan, K.S. (1984). Thermal radiation hazards from hydrocarbon pool fires. Prog. Energy Combust. Sci. 10(1), 59-80.", "Tabela 4 e Tabela 2, correla\u00e7\u00e3o de fuligem"),
    ("9", "Raj, P.K. (2005). Large LNG fire thermal radiation, modeling issues and hazard criteria revisited. Proc. Safety Prog. 24(3), 192-202.", "Tabela 1, as cinco po\u00e7as de GNL"),
    ("10", "Wang, C.J., Wen, J.X., Chen, Z.B. (2014). Simulation of large-scale LNG pool fires using FireFOAM. Comb. Sci. Tech. 186(10-11).", "Comprimento de Montoir, Figura 4"),
    ("11", "Thomas, P.H. (1963). The size of flames from natural fires. 9th Symposium on Combustion, 844-859.", "Comprimento de chama de po\u00e7a"),
    ("12", "Chamberlain, G.A. (1987). Developments in design methods for predicting thermal radiation from flares. Chem. Eng. Res. Des. 65, 299-309.", "Comprimento de chama-jato, Richardson"),
    ("13", "Pietersen, C.M., Huerta, S.C. (1985). Analysis of the LPG incident in San Juan Ixhuatepec. TNO 85-0222.", "Transmissividade atmosf\u00e9rica"),
    ("14", "Eisenberg, N.A., Lynch, C.J., Breeding, R.J. (1975). Vulnerability model. US Coast Guard CG-D-136-75.", "Probit de fatalidade por dose"),
    ("15", "Abramowitz, M., Stegun, I.A. (1964). Handbook of Mathematical Functions, eq. 7.1.26.", "Fun\u00e7\u00e3o erro"),
]
refs_disp = [
    ("16", "Barad, M.L. (1958). Project Prairie Grass. AFCRL-TR-58-235.", "Corridas 7, 11, 22, 29, 35"),
    ("17", "Koopman, R.P. et al. (1982). Burro series data report. LLNL UCID-19075.", "Burro 3, 5, 6, 7, 8, 9"),
    ("18", "Goldwire, H.C. et al. (1983). Coyote series data report. LLNL UCID-19953.", "Coyote 3 e 5"),
    ("19", "Goldwire, H.C. et al. (1985). Desert Tortoise series data report. LLNL UCID-20562.", "Ensaio 4"),
    ("20", "Brown, T.C. et al. (1990). Falcon series data report. GRI-89/0138.", "Falcon 1, 3, 4"),
    ("21", "Puttock, J.S., Blackmore, D.R., Colenbrander, G.W. (1982). Field experiments on dense gas dispersion. J. Haz. Mat. 6, 13-41.", "Maplin Sands 27"),
    ("22", "Roebuck, B. (1983). Thorney Island trials, spill 008. HSE Sheffield.", "Thorney Island 8"),
    ("23", "Biltoft, C.A. (2001). Mock Urban Setting Test. DPG WDTC-FR-01-121. Com Yee e Biltoft (2004).", "MUST 11, 120 obst\u00e1culos"),
    ("24", "Hanna, S.R., Britter, R., Argenta, E., Chang, J. (2012). The Jack Rabbit chlorine release experiments. J. Haz. Mat. 213-214, 406-412.", "Jack Rabbit I, 7"),
    ("25", "Mazzola, T., Hanna, S., Chang, J. et al. (2021). Comparisons of 17 dense gas models to Jack Rabbit II. Atmos. Env. 244, 117887.", "Jack Rabbit II 1 e 7; refer\u00eancia PHAST"),
    ("26", "Witlox, H.W.M., Harper, M., Oke, A. (2012). Phast validation against BP DF1 CO\u2082 experiments. IChemE Hazards XXIII 158; estendido em Witlox et al. (2014).", "Spadeadam DF1 Teste 5"),
    ("27", "Gant, S.E., Kelsey, A., McNally, K., Witlox, H.W.M., Bilio, M. (2014). Multi-phase atmospheric dispersion models for CCS. JLPPI 32, 286-298.", "INERIS CO2PipeHaz 6 mm"),
    ("28", "Gant, S.E., Ivings, M.J. (2005). CFD modelling of low pressure jets for area classification. HSL/2005/13.", "Jato s\u00f4nico de metano"),
    ("29", "Schefer, R.W., Houf, W.G., Williams, T.C. (2008). Small-scale unintended releases of hydrogen. Int. J. Hydrogen Energy 33(21), 6373-6384.", "H\u2082 a 207 bar por 1,91 mm"),
    ("30", "Mack, A., Spruijt, M.P.N. (2013). Validation of OpenFOAM for heavy gas dispersion. J. Haz. Mat. 250-251, 1-14.", "DAT632, SF\u2086 sobre rampa"),
    ("31", "Hansen, O.R., Gavelli, F., Ichard, M., Davis, S.G. (2010). Validation of FLACS for LNG vapor dispersion. JLPPI 23, 857-877.", "Coortes FLACS sem obstru\u00e7\u00e3o e Falcon"),
    ("32", "Hanna, S.R., Hansen, O.R., Dharmavaram, S. (2004). FLACS CFD performance with Kit Fox, MUST, Prairie Grass, EMU. Atmos. Env. 38, 4675-4687.", "Coortes Prairie Grass e MUST; Kit Fox"),
    ("33", "Chang, J.C., Hanna, S.R. (2004). Air quality model performance evaluation. Met. Atmos. Phys. 87, 167-196.", "MRB, RMSE, NMSE, FAC2, MG, VG"),
    ("34", "Vu, T.L. (2019). On numerical modelling of atmospheric gas dispersion using CFD approach. Tese, NTU Singapura.", "Sc\u209c = 0,15, malha, precursor de CLA"),
    ("35", "Briggs, G.A. (1973). Diffusion estimation for small emissions. NOAA ATDL No. 79.", "Ajustes de \u03c3y e \u03c3z"),
    ("36", "Slade, D.H., ed. (1968). Meteorology and Atomic Energy. USAEC TID-24190.", "Coeficientes de sopro instant\u00e2neo"),
    ("37", "Birch, A.D., Brown, D.R., Dodson, M.G., Swaffield, F. (1984). High pressure jets of natural gas. Comb. Sci. Tech. 36, 249-261.", "Fonte expandida em jato s\u00f4nico"),
    ("38", "IOGP (2018). Fire and explosion protection: detection systems. Relat\u00f3rio 434-01.", "Tabela de detectores"),
]


def ref_slide(titulo, sub, itens, n):
    s = slide()
    text(s, 0.85, 0.55, 11.6, 0.3, "Fontes", size=10.5, font=F_MONO,
         color=ACCENT, caps=True, tracking=1.4)
    text(s, 0.85, 0.9, 11.6, 0.5, titulo, size=27, font=F_TITLE, bold=True, color=INK)
    text(s, 0.85, 1.52, 11.6, 0.3, sub, size=11, color=MUTED)
    box(s, 0.85, 1.88, 11.63, 0.012, fill=RULE)
    y = 2.16
    col_w, gap = 5.7, 0.23
    per_col = (len(itens) + 1) // 2
    for i, (num, cit, uso) in enumerate(itens):
        col = 0 if i < per_col else 1
        row = i if i < per_col else i - per_col
        x = 0.85 + col * (col_w + gap)
        ty = y + row * 0.58
        text(s, x, ty, 0.34, 0.25, num, size=9.5, font=F_MONO, color=FAINT,
             align=PP_ALIGN.RIGHT)
        text(s, x + 0.45, ty - 0.015, col_w - 0.5, 0.35, cit, size=8.5, color=INK, spacing=1.18)
        text(s, x + 0.45, ty + 0.36, col_w - 0.5, 0.18, uso, size=7.6, font=F_MONO,
             color=FAINT, spacing=1.1)
    footer(s, n)


ref_slide("Literatura: inc\u00eandio",
          "Experimentos prim\u00e1rios, modelos, correla\u00e7\u00f5es e compila\u00e7\u00f5es de dados.",
          refs_fogo, 19)
ref_slide("Literatura: dispers\u00e3o, experimentos",
          "Ensaios de campo e de t\u00fanel de vento codificados na su\u00edte.",
          refs_disp[:15], 20)
ref_slide("Literatura: dispers\u00e3o, refer\u00eancia e m\u00e9todo",
          "Desempenho publicado dos c\u00f3digos comerciais, o m\u00e9todo estat\u00edstico e as correla\u00e7\u00f5es do motor.",
          refs_disp[15:], 21)

# ════════════════════════════════════════════════════════════════════
# 21. Fecho
# ════════════════════════════════════════════════════════════════════
s = slide(INK)
box(s, 0, 0, W, 0.055, fill=SIGNAL)
text(s, 0.85, 1.5, 11.6, 0.3, "O que este doss\u00ea sustenta", size=11, font=F_MONO,
     color=RGBColor(0x58, 0xBD, 0xB1), caps=True, tracking=1.6)
text(s, 0.85, 2.15, 11.0, 2.0,
     ["47 experimentos de campo,", "271 asser\u00e7\u00f5es, 38 fontes."],
     size=44, font=F_TITLE, bold=True, color=RGBColor(0xF2, 0xF5, 0xF4), spacing=1.1)
box(s, 0.85, 4.5, 11.63, 0.012, fill=RGBColor(0x2C, 0x40, 0x42))
text(s, 0.85, 4.85, 11.0, 1.4,
     ["Todo n\u00famero rastreia at\u00e9 um arquivo de benchmark no reposit\u00f3rio, e todo arquivo carrega a cita\u00e7\u00e3o do experimento que codifica.",
      "As falhas est\u00e3o aqui com o mesmo destaque dos acertos, e o envelope de aplica\u00e7\u00e3o est\u00e1 escrito."],
     size=15, color=RGBColor(0xB9, 0xC7, 0xC7), spacing=1.45, space_after=10)
text(s, 0.85, 6.7, 11.63, 0.3, "DisperSim 3D   \u00b7   28/08/2026",
     size=10, font=F_MONO, color=RGBColor(0x6C, 0x81, 0x82), caps=True, tracking=1.3)

prs.save(OUT)
print("OK:", OUT, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
